import gspread
from oauth2client.service_account import ServiceAccountCredentials

import qrcode
import os

from pptx import Presentation

# ===== GOOGLE SHEETS ULANISH =====

scope = [
    "https://spreadsheets.google.com/feeds",
    "https://www.googleapis.com/auth/drive"
]

creds = ServiceAccountCredentials.from_json_keyfile_name(
    "generator/pure-anthem-496207-h7-0b8d6c1260f9.json",
    scope
)

client = gspread.authorize(creds)

sheet = client.open("sertfikat").sheet1
data = sheet.get_all_records()

# ===== PAPKALAR =====

output_folder = "../pdf"
temp_folder = "temp"

os.makedirs(output_folder, exist_ok=True)
os.makedirs(temp_folder, exist_ok=True)

# ===== HAR BIR TALABA =====

for row in data[147:148]:

    cert_id = row["id"]

    name_uz = row["name_uz"]
    name_en = row["name_en"]

    date_uz = row["date_uz"]
    date_en = row["date_en"]

    hours = row["hours"]

    course_uz = row["course_uz"]
    course_en = row["course_en"]

    issue_date = row["issue_date"]

    status = row["status"]
    work_type = row["work_type"]

    # ===== STATUS =====

    if str(status).lower() != "otdi":
        continue

    # ===== TRAINING TYPE =====

    if work_type == "To'lov shartnoma":

        training_type = "qayta tayyorlash"
        training_type_en = "retraining"

    else:

        training_type = "malaka oshirish"
        training_type_en = "skill development"

    # ===== QR =====

    qr_link = f"https://olmaliq-kmk.github.io/Engineering_school/?id={cert_id}"

    qr = qrcode.make(qr_link)

    qr_path = f"{temp_folder}/{cert_id}.png"

    qr.save(qr_path)

    # ===== PPTX TEMPLATE =====

    prs = Presentation("generator/shablon.pptx")

    replacements = {

        "{{ID}}": str(cert_id),

        "{{NAME_UZ}}": str(name_uz),
        "{{NAME_EN}}": str(name_en),

        "{{DATE_UZ}}": str(date_uz),
        "{{DATE_EN}}": str(date_en),

        "{{HOURS}}": str(hours),

        "{{COURSE_UZ}}": str(course_uz),
        "{{COURSE_EN}}": str(course_en),

        "{{ISSUE_DATE}}": str(issue_date),

        "{{TRAINING_TYPE}}": str(training_type),
        "{{TRAINING_TYPE_EN}}": str(training_type_en),

    }

 

# ===== TEXT ALMASHTIRISH =====

for slide in prs.slides:

    for shape in slide.shapes:

        if not hasattr(shape, "text"):
            continue

        for key, value in replacements.items():

            if key in shape.text:

                shape.text = shape.text.replace(key, value)

    # ===== QR QO'SHISH =====

    for slide in prs.slides:

        slide.shapes.add_picture(
            qr_path,
            left=5400000,
            top=5100000,
            width=900000,
            height=900000
        )

    # ===== PPTX SAQLASH =====

    pptx_path = f"{temp_folder}/{cert_id}.pptx"

    prs.save(pptx_path)

    print(f"✅ Tayyor: {cert_id}")

print("🔥 BARCHA SERTIFIKATLAR TAYYOR")